import os
import fitz  # PyMuPDF
from pptx import Presentation
from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse, HttpResponse
from django.views.decorators.http import require_POST
from django.contrib.auth.decorators import login_required
from django.core.exceptions import PermissionDenied
from .forms import ChatForm
from .models import Chat, Message
from .utils import require_auth, generate_auth_token
import google.generativeai as genai
import dotenv
import json
import uuid
from django.template.loader import render_to_string
import markdown2
import tempfile
import pdfkit
from xhtml2pdf import pisa
from io import BytesIO

dotenv.load_dotenv()
API = os.getenv("GOOGLE_API_KEY")

genai.configure(api_key=API)
model = genai.GenerativeModel("gemini-2.5-flash-preview-05-20")

def extract_text_from_pdf(file):
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_pptx(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def index(request):
    if request.user.is_authenticated:
        return redirect('chat')
    return render(request, 'index.html')

@login_required
@require_auth
def chatbot_view(request, chat_id=None):
    # Get all chats for the current user
    chats = Chat.objects.filter(user=request.user).order_by('-updated_at')
    current_chat_title = None
    for chat in chats:
        if chat.id == chat_id:
            current_chat_title = chat.title
    
    # Generate auth token for the frontend
    auth_token = generate_auth_token()
    
    # If no chat_id is provided and no chats exist, just show the empty state
    if chat_id is None and not chats.exists():
        return render(request, 'chatbot.html', {
            'chats': chats,
            'show_empty_state': True,
            'auth_token': auth_token
        })
    
    # If no chat_id is provided but chats exist, redirect to the most recent chat
    if chat_id is None:
        return redirect('chat', chat_id=chats.first().id)
    
    # Convert chat_id to UUID if it's a string
    if isinstance(chat_id, str):
        try:
            chat_id = uuid.UUID(chat_id)
        except ValueError:
            raise PermissionDenied("Invalid chat ID.")
    
    # Get the chat and verify ownership
    chat = get_object_or_404(Chat, id=chat_id)
    if chat.user != request.user:
        raise PermissionDenied("You don't have permission to access this chat.")
    
    messages = chat.messages.all()
    
    if request.method == 'POST':
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            user_message = request.POST.get('message', '')
            print(user_message)
            # Get list of uploaded files (up to 3)
            uploaded_files = request.FILES.getlist('file')[:3]
            print(uploaded_files)
            
            combined_file_text = ""
            uploaded_file_names = []

            if uploaded_files:
                for uploaded_file in uploaded_files:
                    uploaded_file_names.append(uploaded_file.name)
                    ext = os.path.splitext(uploaded_file.name)[-1].lower()
                    try:
                        if ext == '.pdf':
                            combined_file_text += f"Document Name: {uploaded_file.name}\n" + extract_text_from_pdf(uploaded_file) + "\n\n"
                        elif ext == '.pptx':
                            combined_file_text += f"Document Name: {uploaded_file.name}\n" + extract_text_from_pptx(uploaded_file) + "\n\n"
                    except Exception as e:
                        print(f"Error processing file {uploaded_file.name}: {e}")

            # Create user message
            Message.objects.create(
                chat=chat,
                sender='You',
                text=user_message,
                attachment=','.join(uploaded_file_names) if uploaded_file_names else None
            )
            
            # Format chat history for context
            history = ""
            for msg in messages:
                if msg.sender == 'You':
                    history += f"You: {msg.text}\n"
                else:
                    history += f"PrepAI: {msg.text}\n"

            # Include exam context in the prompt if available
            exam_context = ""
            if chat.exam_name:
                exam_context = f"\nExam Context:\n- Exam Name: {chat.exam_name}\n"
                if chat.exam_details:
                    exam_context += f"- Exam Details: {chat.exam_details}\n"

            prompt = (
                "You are PrepAI, an intelligent and friendly assistant designed to help students prepare for exams. "
                "You are made by a company called PrepAI. Only reveal this if asked."
                "You are knowledgeable, supportive, and always focused on helping students achieve their academic goals.\n\n"
                + exam_context
                + (f"Here is the previous conversation for context:\n{history}\n" if history else "")
                + "## Your Capabilities:\n"
                "- Analyze study materials (PDFs, PPTX, notes) and extract key information\n"
                "- Generate practice questions (MCQs, short answers, long answers)\n"
                "- Create study summaries and highlight important concepts\n"
                "- Provide personalized study tips and strategies\n"
                "- Suggest study schedules and revision plans\n"
                "- Answer subject-specific questions\n"
                "- Help with exam preparation strategies\n\n"
                "## Response Guidelines:\n"
                "- If the user asks about a specific topic, then explain the topic to the user topic wise and ask the user for a follow-up question\n"
                "- If the user is asking about the document, then divide the document into topics and explain it to the user topic wise and ask the user for a follow-up question\n"
                "- Always be helpful, encouraging, and student-focused\n"
                "- Keep explanations clear and avoid unnecessary jargon\n"
                "- Provide actionable advice and specific examples\n"
                "- Don't mention disclaimers like 'I've analyzed' or 'Based on the document'\n"
                "- If asked about non-academic topics, politely redirect to exam preparation\n"
                "- If greeting (like 'Hi', 'Hello'), respond warmly and ask how you can help with exam prep\n"
                "- If given random/off-topic information, acknowledge briefly and guide back to studying\n\n"
            )

            if combined_file_text.strip():
                prompt += (
                    f"### Document Analysis Task:\n"
                    f"- Document Text: {combined_file_text}\n"
                    f"- Student's Message: {user_message}\n\n"
                    "Your job is to analyze the provided material and assist the student by:\n"
                    "- Summarizing the important points relevant to the exam\n"
                    "- Highlighting key concepts likely to be asked\n"
                    "- Generating potential questions (MCQs, short/long answers) based on the content\n"
                    "- Giving personalized study tips based on the material\n"
                    "- Identifying gaps or suggesting additional topics to review\n"
                    "- If this is the first message from the user, then divide the document into topics and explain the topics to the user topic wise and ask the user for a follow-up question\n\n"
                )
            else:
                prompt += (
                    f"### Conversation Task:\n"
                    f"- Student's Message: {user_message}\n\n"
                    "Since no document was provided, focus on:\n"
                    "- Answering subject-specific questions with detailed explanations\n"
                    "- Providing study tips and strategies based on their exam context\n"
                    "- Creating practice questions if requested\n"
                    "- Offering general exam preparation guidance\n"
                    "- If this is a greeting, welcome them and ask about their exam preparation needs\n"
                    "- If they share off-topic information, acknowledge it and redirect to exam prep\n\n"
                )

            prompt += (
                "## Response Requirements:\n"
                "- Provide clear, concise information directly related to exam preparation\n"
                "- List key points or concepts when analyzing documents\n"
                "- Generate questions or study tips aligned with the exam context\n"
                "- Maintain an encouraging and supportive tone throughout\n"
                "- Always focus on helping the student succeed in their exams\n\n"
                "Always keep your explanations student-friendly, focused, and aligned with the exam goal.\n\n"
                "Now respond helpfully to assist the student with their exam preparation!"
            )

            response = model.generate_content(prompt)
            response_text = response.text
            
            # Create LLM message
            Message.objects.create(
                chat=chat,
                sender='PrepAI',
                text=response_text
            )
            
            # Update chat title if it's the first message
            if messages.count() == 0:
                chat.title = user_message[:50] + "..." if len(user_message) > 50 else user_message
                chat.save()
            
            # Return updated messages for frontend
            messages = chat.messages.all()
            return JsonResponse({
                'response': response_text,
                'messages': list(messages.values('sender', 'text', 'attachment', 'created_at', 'id'))
            })

    else:
        form = ChatForm()
        for msg in messages:
            if msg.attachment:
                msg.attachments_list = [att.strip() for att in msg.attachment.split(',')]
            else:
                msg.attachments_list = []
    
    return render(request, 'chatbot.html', {
        'form': form,
        'messages': messages,
        'chats': chats,
        'current_chat_id': str(chat.id),
        'auth_token': auth_token,
        'current_chat_title': current_chat_title
    })

@login_required
@require_auth
@require_POST
def create_chat(request):
    exam_name = request.POST.get('exam_name')
    exam_details = request.POST.get('exam_details', '')
    if 'testuser' in str(request.user):
        if len(Chat.objects.filter(user=request.user)) >= 1:
            return JsonResponse({'success': False, 'error': 'Only one chat is allowed for anonymous users'}, status=400)

    if not exam_name:
        return JsonResponse({'success': False, 'error': 'Exam name is required'}, status=400)
    
    chat = Chat.objects.create(
        title=exam_name,
        user=request.user,
        exam_name=exam_name,
        exam_details=exam_details
    )
    return JsonResponse({'success': True, 'chat_id': str(chat.id)})

@login_required
@require_auth
@require_POST
def delete_chat(request, chat_id):
    # Convert chat_id to UUID if it's a string
    if isinstance(chat_id, str):
        try:
            chat_id = uuid.UUID(chat_id)
        except ValueError:
            raise PermissionDenied("Invalid chat ID.")
    chat = get_object_or_404(Chat, id=chat_id)
    if chat.user != request.user:
        raise PermissionDenied("You don't have permission to delete this chat.")
    chat.delete()
    return JsonResponse({'success': True})

@login_required
def download_message_pdf(request, message_id):
    try:
        message = Message.objects.get(id=message_id, chat__user=request.user)
        html_content = markdown2.markdown(message.text, safe_mode='escape', extras=['fenced-code-blocks', 'tables', 'code-friendly'])

        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                body {{ font-family: Arial, sans-serif; padding: 20px; }}
            </style>
        </head>
        <body>
            {html_content}
        </body>
        </html>
        """

        # Create PDF in memory
        pdf_buffer = BytesIO()
        pisa_status = pisa.CreatePDF(full_html, dest=pdf_buffer)

        if pisa_status.err:
            return HttpResponse("Error generating PDF", status=500)

        # Return PDF as response
        response = HttpResponse(pdf_buffer.getvalue(), content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="PrepAI_Notes.pdf"'
        return response
        
    except Message.DoesNotExist:
        return HttpResponse('Message not found', status=404)
    except Exception as e:
        return HttpResponse(f'Error generating PDF: {str(e)}', status=500)
